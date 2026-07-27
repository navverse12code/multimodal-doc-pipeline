from docx import Document
from pptx import Presentation
from pathlib import Path

docs_dir = Path("sample_docs")
docs_dir.mkdir(exist_ok=True)

# 1. Create DOCX
doc = Document()
doc.add_heading('Vendor Agreement 2026', 0)
doc.add_paragraph('This document outlines the security policies and server license pricing for Acme Corp.')
table = doc.add_table(rows=1, cols=3)
hdr_cells = table.rows[0].cells
hdr_cells[0].text = 'Item'
hdr_cells[1].text = 'Quantity'
hdr_cells[2].text = 'Price'
row_cells = table.add_row().cells
row_cells[0].text = 'Database Storage'
row_cells[1].text = '50 TB'
row_cells[2].text = '$4,200'
doc.save(docs_dir / 'vendor_agreement.docx')

# 2. Create PPTX
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "Q3 Product Roadmap"
slide.placeholders[1].text = "Key Objectives:\n1. Multi-modal AI Search\n2. SQLite FTS5 Integration"
prs.save(docs_dir / 'product_roadmap.pptx')

print("Sample DOCX and PPTX created successfully!")
