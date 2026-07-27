import logging
from pptx import Presentation
from parsers.base_parser import BaseParser

logger = logging.getLogger(__name__)

class PPTXParser(BaseParser):
    def parse(self, file_path: str):
        slides = []
        try:
            prs = Presentation(file_path)
            for idx, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                slide_text.append(text)
                                
                combined_slide = "\n".join(slide_text)
                if combined_slide.strip():
                    slides.append((idx + 1, combined_slide))

        except Exception as e:
            logger.error(f"Error parsing PPTX {file_path}: {e}")

        return slides
